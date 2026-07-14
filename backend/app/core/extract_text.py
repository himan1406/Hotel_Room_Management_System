import io

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def extract_text_from_bytes(content: bytes, filename: str) -> str:
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    if ext == "pdf":
        return _extract_pdf(content)
    elif ext in ("docx", "doc"):
        return _extract_docx(content)
    elif ext == "xlsx":
        return _extract_xlsx(content)
    elif ext == "pptx":
        return _extract_pptx(content)
    elif ext in ("txt", "csv", "json", "xml", "md"):
        return content.decode("utf-8", errors="replace")
    else:
        return ""


def _extract_pdf(content: bytes) -> str:
    if fitz is None:
        return ""
    texts = []
    doc = fitz.open(stream=content, filetype="pdf")
    for page in doc:
        texts.append(page.get_text())
    doc.close()
    return "\n".join(texts)


def _extract_docx(content: bytes) -> str:
    if DocxDocument is None:
        return ""
    doc = DocxDocument(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_xlsx(content: bytes) -> str:
    if openpyxl is None:
        return ""
    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    texts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None]
            if cells:
                texts.append(" ".join(cells))
    wb.close()
    return "\n".join(texts)


def _extract_pptx(content: bytes) -> str:
    if Presentation is None:
        return ""
    prs = Presentation(io.BytesIO(content))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)
